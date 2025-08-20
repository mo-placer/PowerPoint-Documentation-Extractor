#!/usr/bin/env python3
"""
PowerPoint Documentation Updates Extractor

This script extracts multiple components from PowerPoint files:
1. Vocabulary words and definitions (blue/turquoise and bold text)
2. Session Goals (bulleted lists after "In today's session, you will")
3. Assessments (numbered lists after instructor evaluation text)
4. Related Careers (bulleted list on the last slide; trigger: "Related Careers" or "Careers")
5. Session Materials (bulleted list; triggers like "Gather the following items", "Gather the following item",
   "Locate the following", or "Locate the ...")
"""

import os
import re
from pathlib import Path
from typing import List, Tuple, Dict
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
except ImportError as e:
    print(f"Required packages not installed: {e}")
    print("Please install required packages:")
    print("pip install python-pptx python-docx")
    exit(1)


class DocumentationExtractor:
    def __init__(self):
        # Define blue/turquoise color ranges (RGB values)
        self.blue_color_ranges = [
            # Standard blues
            (0, 0, 200, 255, 255, 255),  # Blue range
            (0, 100, 150, 100, 255, 255),  # Turquoise range
            (0, 150, 200, 150, 255, 255),  # Light blue range
            (30, 144, 255, 70, 180, 255),  # Dodger blue range
            (0, 191, 255, 50, 220, 255),  # Deep sky blue range
            (64, 224, 208, 100, 255, 230),  # Turquoise range
        ]
        
        # Common headings/placeholders that are NOT vocabulary even if bold/blue
        self.non_vocab_terms_normalized = {
            'direct instruction', 'procedure', 'direct instruction procedure', 'procedure safety',
            'direct instruction, procedure', 'direct instruction , procedure',  # Handle comma variations
            'vocab definition', 'module guide', 'assessment', 'goals', 'gather', 'clean up',
            'summary', 'fa kc', 'activity lead in', 'post test', 'intro', 'worksheet',
            'careers', 'related careers', 'materials', 'session materials',
            'gather the following items', 'gather the following item',
            'locate the following', 'locate the',
            'insert slide type', 'insert slide layout', 'insert media description'
        }

    def _normalize_phrase(self, text: str) -> str:
        return re.sub(r'[^a-z0-9]+', ' ', text.lower()).strip()

    def should_exclude_vocab(self, text: str) -> bool:
        if not text:
            return True
        
        # Clean up the text for analysis
        clean_text = text.strip()
        normalized = self._normalize_phrase(clean_text)
        
        # Exclude very short tokens and lone letters
        if len(normalized) <= 2:
            return True
            
        # Exclude single letters or very short meaningless text
        if len(clean_text) <= 2:
            return True
        
        # AGGRESSIVE exclusion patterns for common slide layout placeholders
        placeholder_patterns = [
            'direct instruction', 'procedure', 'worksheet', 'vocab definition',
            'activity lead in', 'gather', 'clean up', 'summary', 'goals',
            'assessment', 'careers', 'related careers', 'materials', 
            'session materials', 'fa kc', 'post test', 'intro',
            'insert slide', 'insert media', 'insert layout', 'objectives'
        ]
        
        # Check if the text contains ANY of these placeholder patterns
        for pattern in placeholder_patterns:
            if pattern in normalized:
                return True
        
        # Check for exact matches in our exclusion set
        if normalized in self.non_vocab_terms_normalized:
            return True
            
        # Exclude anything that starts with common placeholder words
        placeholder_starts = ['insert', 'post', 'direct', 'procedure', 'vocab', 'gather', 'clean']
        for start_word in placeholder_starts:
            if normalized.startswith(start_word + ' '):
                return True
        
        # Exclude anything that ends with common placeholder indicators
        if normalized.endswith(' worksheet') or normalized.endswith(' definition'):
            return True
            
        # Exclude comma-separated combinations of placeholder terms
        if ',' in clean_text:
            parts = [part.strip() for part in clean_text.split(',')]
            all_parts_are_placeholders = True
            for part in parts:
                part_normalized = self._normalize_phrase(part)
                is_placeholder = False
                for pattern in placeholder_patterns:
                    if pattern in part_normalized or part_normalized in pattern:
                        is_placeholder = True
                        break
                if not is_placeholder:
                    all_parts_are_placeholders = False
                    break
            if all_parts_are_placeholders:
                return True
        
        return False
    
    def resolve_theme_color_to_rgb(self, presentation, theme_color_idx, brightness=0):
        """Resolve a theme color index to RGB values."""
        try:
            theme = presentation.slide_master.theme_part.theme
            color_scheme = theme.theme_elements.clrScheme
            
            # Map theme color indices to color scheme elements
            theme_color_map = {
                1: color_scheme.dk1_color,  # Dark 1
                2: color_scheme.lt1_color,  # Light 1
                3: color_scheme.dk2_color,  # Dark 2
                4: color_scheme.lt2_color,  # Light 2
                5: color_scheme.accent1_color,  # Accent 1
                6: color_scheme.accent2_color,  # Accent 2
                7: color_scheme.accent3_color,  # Accent 3
                8: color_scheme.accent4_color,  # Accent 4
                9: color_scheme.accent5_color,  # Accent 5
                10: color_scheme.accent6_color,  # Accent 6
                11: color_scheme.hlink_color,   # Hyperlink
                12: color_scheme.folHlink_color  # Followed Hyperlink
            }
            
            if theme_color_idx in theme_color_map:
                theme_color = theme_color_map[theme_color_idx]
                if hasattr(theme_color, 'rgb'):
                    return theme_color.rgb
                    
        except Exception as e:
            pass
        return None

    def get_effective_color(self, run, paragraph, presentation):
        """Get the effective color of a text run using layered approach."""
        # Layer 1: Check run-level color
        if run.font.color:
            # Explicit RGB color
            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                return run.font.color.rgb
            
            # Theme color
            if hasattr(run.font.color, 'theme_color') and run.font.color.theme_color:
                brightness = getattr(run.font.color, 'brightness', 0)
                rgb = self.resolve_theme_color_to_rgb(presentation, run.font.color.theme_color, brightness)
                if rgb:
                    return rgb
        
        # Layer 2: Check paragraph-level color
        if hasattr(paragraph, 'font') and paragraph.font.color:
            if hasattr(paragraph.font.color, 'rgb') and paragraph.font.color.rgb:
                return paragraph.font.color.rgb
            if hasattr(paragraph.font.color, 'theme_color') and paragraph.font.color.theme_color:
                brightness = getattr(paragraph.font.color, 'brightness', 0)
                rgb = self.resolve_theme_color_to_rgb(presentation, paragraph.font.color.theme_color, brightness)
                if rgb:
                    return rgb
        
        # Layer 3: Fall back to theme defaults (Text1/Text2)
        # Assume Text1 (theme color 1) for now
        return self.resolve_theme_color_to_rgb(presentation, 1)

    def is_blue_color(self, rgb) -> bool:
        """Check if RGB values fall within blue/turquoise ranges."""
        if rgb is None:
            return False
        
        try:
            # RGBColor uses different attribute names
            r = rgb.red if hasattr(rgb, 'red') else (rgb[0] if isinstance(rgb, tuple) else rgb.r)
            g = rgb.green if hasattr(rgb, 'green') else (rgb[1] if isinstance(rgb, tuple) else rgb.g)
            b = rgb.blue if hasattr(rgb, 'blue') else (rgb[2] if isinstance(rgb, tuple) else rgb.b)
            
            # Check against all blue color ranges
            for r_min, g_min, b_min, r_max, g_max, b_max in self.blue_color_ranges:
                if (r_min <= r <= r_max and 
                    g_min <= g <= g_max and 
                    b_min <= b <= b_max):
                    return True
                    
            # Additional specific blue checks
            # Exclude white, near-white, and gray colors
            if r > 200 and g > 200 and b > 200:
                return False
            
            # Exclude pure gray colors
            if abs(r - g) < 10 and abs(g - b) < 10 and abs(r - b) < 10:
                return False
                
            # Check for common PowerPoint blues (must be significantly more blue)
            if (b > r + 50 and b > g + 50 and b > 120):  # Strong blue dominance
                return True
                
            # Check for turquoise (high green and blue, low red)
            if (g > 150 and b > 150 and r < 80):
                return True
                
            # Check for medium blues with good blue component
            if (b > 180 and b > r + 30 and b > g + 30):
                return True
                
        except:
            pass
            
        return False
    
    def extract_vocabulary_from_slide(self, slide, presentation, debug=False) -> List[Tuple[str, str]]:
        """Extract vocabulary words and definitions from a single slide."""
        vocabulary_items = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                
                for paragraph in text_frame.paragraphs:
                    vocab_word = None
                    definition_parts = []
                    current_paragraph_text = paragraph.text.strip()
                    
                    if not current_paragraph_text:
                        continue
                    
                    if debug:
                        print(f"    Analyzing paragraph: {current_paragraph_text[:100]}...")
                    
                    for run in paragraph.runs:
                        text = run.text.strip()
                        if not text:
                            continue
                        
                        # Get effective color using layered approach
                        effective_rgb = self.get_effective_color(run, paragraph, presentation)
                        is_blue = self.is_blue_color(effective_rgb)
                        is_bold = run.font.bold
                        
                        if debug:
                            color_info = "None"
                            if effective_rgb:
                                try:
                                    r = effective_rgb.red if hasattr(effective_rgb, 'red') else effective_rgb[0]
                                    g = effective_rgb.green if hasattr(effective_rgb, 'green') else effective_rgb[1] 
                                    b = effective_rgb.blue if hasattr(effective_rgb, 'blue') else effective_rgb[2]
                                    color_info = f"RGB({r}, {g}, {b})"
                                except:
                                    color_info = f"RGB(error)"
                            
                            theme_color_info = "None"
                            if run.font.color and hasattr(run.font.color, 'theme_color'):
                                theme_color_info = f"Theme({run.font.color.theme_color})"
                            
                            print(f"      Run: '{text}' | Bold: {is_bold} | Blue: {is_blue} | Color: {color_info} | Theme: {theme_color_info}")
                        
                        if is_blue and is_bold and not vocab_word and not self.should_exclude_vocab(text):
                            # This is likely a vocabulary word
                            vocab_word = text
                            if debug:
                                print(f"        -> Found vocabulary word: {vocab_word}")
                        elif vocab_word and text:
                            # This might be part of the definition
                            definition_parts.append(text)
                    
                    # Alternative approach: look for bold text that might be vocabulary
                    if not vocab_word:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            if text and run.font.bold and not self.should_exclude_vocab(text):
                                # Check if this looks like a vocabulary word (single word or short phrase)
                                if (1 < len(text.split()) <= 3 and 
                                    not any(char in text for char in '.,!?;[]()') and
                                    not text.isdigit() and
                                    len(text) > 2):  # Must be more than 2 characters
                                    vocab_word = text
                                    if debug:
                                        print(f"        -> Found potential vocabulary (bold): {vocab_word}")
                                    break
                    
                    # If we found a vocabulary word, extract the definition
                    if vocab_word:
                        # Get the rest of the paragraph as definition
                        full_text = paragraph.text
                        
                        # Try to extract definition after the vocabulary word
                        vocab_pattern = re.escape(vocab_word)
                        match = re.search(f'{vocab_pattern}\\s*[:-]?\\s*(.+)', full_text, re.IGNORECASE)
                        
                        if match:
                            definition = match.group(1).strip()
                        else:
                            definition = ' '.join(definition_parts).strip()
                            # If no definition parts, try to get text after the vocab word
                            if not definition:
                                text_after_vocab = full_text[full_text.lower().find(vocab_word.lower()) + len(vocab_word):].strip()
                                if text_after_vocab.startswith((':', '-', '–', '—')):
                                    definition = text_after_vocab[1:].strip()
                                else:
                                    definition = text_after_vocab
                        
                        if definition and len(definition) > 3:  # Ensure we have a meaningful definition
                            vocabulary_items.append((vocab_word, definition))
                            if debug:
                                print(f"        -> Added: {vocab_word} = {definition}")
        
        return vocabulary_items
    
    def extract_session_goals(self, slide, debug=False) -> List[str]:
        """Extract session goals from a slide. Captures everything after trigger text within the slide."""
        goals: List[str] = []
        trigger_patterns = [
            r"In today's session,?\s*you will:?",
            r"In this session,?\s*you will:?",
            r"Today,?\s*you will:?",
            r"^GOALS$",  # Simple GOALS heading (exact match)
            r"^Goals$",  # Goals heading (title case, exact match)
            r"Session Goals",  # Session Goals heading
            r"Learning Goals",  # Learning Goals heading
            r"^Objectives$"  # Objectives heading (exact match)
        ]

        # Pass 1: detect if any shape on this slide contains a goals trigger
        slide_has_goals = False
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                full_text = shape.text_frame.text
                if any(re.search(pat, full_text, re.IGNORECASE) for pat in trigger_patterns):
                    slide_has_goals = True
                    if debug:
                        print("    Goals trigger found on slide")
                    break

        if not slide_has_goals:
            return goals

        # Pass 2: Extract all content after trigger patterns from all shapes
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.text_frame.text:
                continue
            tf = shape.text_frame
            full_text = tf.text

            # Check if this shape contains a trigger pattern
            trigger_found = False
            trigger_end_pos = 0
            for pattern in trigger_patterns:
                match = re.search(pattern, full_text, re.IGNORECASE)
                if match:
                    trigger_found = True
                    trigger_end_pos = match.end()
                    break

            if trigger_found:
                # Extract everything after the trigger pattern
                content_after_trigger = full_text[trigger_end_pos:].strip()
                if debug:
                    print(f"      Content after trigger: {content_after_trigger[:100]}...")
                
                # Split into lines and extract ALL content after trigger
                lines = content_after_trigger.split('\n')
                for line in lines:
                    line = line.strip()
                    if line and len(line) > 3:  # Skip very short lines
                        # Skip any remaining trigger text fragments and instructional text
                        line_lower = line.lower()
                        if (any(trigger_fragment in line_lower for trigger_fragment in [
                            'in today\'s session', 'in this session', 'today, you will', 'you will:',
                            'goals:', 'objectives:', 'session goals', 'learning goals'
                        ]) or 
                        # Skip lines that are just trigger words
                        line_lower.strip() in ['goals', 'objectives', 'session goals', 'learning goals'] or
                        # Skip long instructional sentences
                        any(instruction in line_lower for instruction in [
                            'now that you have completed', 'module activity sheet', 'student portfolio',
                            'turn on your call light', 'review the following skills', 'waiting for your instructor'
                        ])):
                            continue
                        
                        # Clean up bullet points and list numbering (preserve quantities)
                        clean_goal = re.sub(r'^\d+[\.\)]\s+', '', line)  # Remove "1. " or "2) " but not "2 "
                        clean_goal = clean_goal.lstrip('•·-*◦').strip()  # Remove bullets
                        if clean_goal and clean_goal not in goals:
                            goals.append(clean_goal)
                            if debug:
                                print(f"        -> Goal: {clean_goal}")
            else:
                # If no trigger in this shape, but slide has goals, check for standalone goal items
                lines = full_text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line and len(line) > 10:  # Reasonable length for a goal
                        # Skip obvious non-goals and trigger text
                        line_lower = line.lower()
                        if (any(skip_term in line_lower for skip_term in [
                            'insert', 'slide layout', 'media description', 'copyright', '©'
                        ]) or
                        any(trigger_fragment in line_lower for trigger_fragment in [
                            'in today\'s session', 'in this session', 'today, you will', 'you will:',
                            'goals:', 'objectives:', 'session goals', 'learning goals'
                        ]) or
                        line_lower.strip() in ['goals', 'objectives', 'session goals', 'learning goals']):
                            continue
                        
                        # If it starts with bullet points or looks like a goal statement, include it
                        if (line.startswith(('•', '·', '-', '*', '◦')) or
                            # Or if it contains typical goal-like language patterns
                            any(word in line_lower for word in [
                                'will', 'explore', 'examine', 'learn', 'understand', 'identify',
                                'analyze', 'demonstrate', 'observe', 'discover', 'explain',
                                'describe', 'investigate', 'construct', 'compare', 'use',
                                'apply', 'evaluate', 'create', 'synthesize', 'capture',
                                'develop', 'practice', 'determine', 'calculate', 'solve',
                                'recognize', 'review', 'begin', 'complete'
                            ])):
                            clean_goal = re.sub(r'^\d+[\.\)]\s+', '', line)  # Remove "1. " or "2) " but not "2 "
                            clean_goal = clean_goal.lstrip('•·-*◦').strip()
                            if clean_goal and clean_goal not in goals:
                                goals.append(clean_goal)
                                if debug:
                                    print(f"        -> Goal (standalone): {clean_goal}")
                
                # Also check paragraph-level for goals that might not be line-separated
                for paragraph in tf.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text and len(para_text) > 10:  # Reasonable length for a goal
                        para_lower = para_text.lower()
                        # Skip trigger lines, obvious non-goals, and trigger text
                        if (any(re.search(pat, para_text, re.IGNORECASE) for pat in trigger_patterns) or
                            any(skip_term in para_lower for skip_term in [
                                'insert', 'slide layout', 'media description', 'copyright', '©'
                            ]) or
                            any(trigger_fragment in para_lower for trigger_fragment in [
                                'in today\'s session', 'in this session', 'today, you will', 'you will:',
                                'goals:', 'objectives:', 'session goals', 'learning goals'
                            ]) or
                            para_lower.strip() in ['goals', 'objectives', 'session goals', 'learning goals']):
                            continue
                        
                        # If it contains goal-like language, include it
                        if (para_text.startswith(('•', '·', '-', '*', '◦')) or
                            any(word in para_lower for word in [
                                'will', 'explore', 'examine', 'learn', 'understand', 'identify',
                                'analyze', 'demonstrate', 'observe', 'discover', 'explain',
                                'describe', 'investigate', 'construct', 'compare', 'use',
                                'apply', 'evaluate', 'create', 'synthesize', 'capture',
                                'develop', 'practice', 'determine', 'calculate', 'solve',
                                'recognize', 'review', 'begin', 'complete'
                            ])):
                            clean_goal = re.sub(r'^\d+[\.\)]\s+', '', para_text)  # Remove "1. " or "2) " but not "2 "
                            clean_goal = clean_goal.lstrip('•·-*◦').strip()
                            if clean_goal and clean_goal not in goals:
                                goals.append(clean_goal)
                                if debug:
                                    print(f"        -> Goal (para standalone): {clean_goal}")

        return goals
    
    def extract_assessments(self, slide, debug=False) -> List[str]:
        """Extract assessment items from a slide. Captures everything after trigger text within the slide."""
        assessments: List[str] = []

        # Expanded trigger matching to handle various phrasing variations
        trigger_patterns = [
            r"your instructor will be evaluating you",
            r"listed onscreen are the specific items",
            r"assessment",  # heading often present on the slide
            r"you will be evaluated on",
            r"evaluation criteria",
            r"instructor will evaluate",
            r"assessment criteria",
            r"performance indicators",
            r"you should be able to"  # Common lead-in for assessment items
        ]

        # Pass 1: detect if the slide is an assessment slide
        slide_has_assessment = False
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                full_text = shape.text_frame.text
                if any(re.search(pat, full_text, re.IGNORECASE) for pat in trigger_patterns):
                    slide_has_assessment = True
                    if debug:
                        print("    Assessment trigger found on slide")
                    break

        if not slide_has_assessment:
            return assessments

        # Pass 2: Extract all content after trigger patterns from all shapes
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.text_frame.text:
                continue
            tf = shape.text_frame
            full_text = tf.text

            # Check if this shape contains a trigger pattern
            trigger_found = False
            trigger_end_pos = 0
            for pattern in trigger_patterns:
                match = re.search(pattern, full_text, re.IGNORECASE)
                if match:
                    trigger_found = True
                    trigger_end_pos = match.end()
                    break

            if trigger_found:
                # Extract everything after the trigger pattern
                content_after_trigger = full_text[trigger_end_pos:].strip()
                if debug:
                    print(f"      Content after trigger: {content_after_trigger[:100]}...")
                
                # Split into lines and extract ALL content after trigger
                lines = content_after_trigger.split('\n')
                for line in lines:
                    line = line.strip()
                    if line and len(line) > 3:  # Skip very short lines
                        # Skip common instructional text that's not an assessment
                        line_lower = line.lower()
                        if (("review the list" in line_lower) or ("before continuing" in line_lower) or
                            # Skip long instructional sentences
                            any(instruction in line_lower for instruction in [
                                'now that you have completed', 'module activity sheet', 'student portfolio',
                                'turn on your call light', 'review the following skills', 'waiting for your instructor',
                                'find the module activity sheet'
                            ])):
                            continue
                        
                        # Clean up bullet points and list numbering (preserve quantities)
                        clean_assessment = re.sub(r'^\d+[\.\)]\s+', '', line)  # Remove "1. " or "2) " but not "2 "
                        clean_assessment = clean_assessment.lstrip('•·-*◦').strip()  # Remove bullets
                        if clean_assessment and clean_assessment not in assessments:
                            assessments.append(clean_assessment)
                            if debug:
                                print(f"        -> Assessment: {clean_assessment}")
            else:
                # If no trigger in this shape, but slide has assessments, check for standalone assessment items
                lines = full_text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line and len(line) > 10:  # Reasonable length for an assessment
                        line_lower = line.lower()
                        # Skip trigger text, instructional text, and obvious non-assessments
                        if (any(re.search(pat, line, re.IGNORECASE) for pat in trigger_patterns) or
                            "review the list" in line_lower or "before continuing" in line_lower or
                            any(skip_term in line_lower for skip_term in [
                                'insert', 'slide layout', 'media description', 'copyright', '©'
                            ]) or
                            # Skip long instructional sentences
                            any(instruction in line_lower for instruction in [
                                'now that you have completed', 'module activity sheet', 'student portfolio',
                                'turn on your call light', 'review the following skills', 'waiting for your instructor',
                                'find the module activity sheet'
                            ])):
                            continue
                        
                        # If it looks like an assessment item, include it
                        if (re.match(r'^\d+\.?\s+', line) or
                            line.startswith(('•', '·', '-', '*', '◦')) or
                            # Or if it contains assessment-like language
                            any(word in line_lower for word in [
                                'describe','explain','identify','analyze','compare','demonstrate',
                                'use','construct','examine','define','name','list','calculate',
                                'determine','give','provide','show','illustrate','evaluate'
                            ])):
                            clean_assessment = re.sub(r'^\d+[\.\)]\s+', '', line)  # Remove "1. " or "2) " but not "2 "
                            clean_assessment = clean_assessment.lstrip('•·-*◦').strip()
                            if clean_assessment and clean_assessment not in assessments:
                                assessments.append(clean_assessment)
                                if debug:
                                    print(f"        -> Assessment (standalone): {clean_assessment}")

        return assessments
    
    def _clean_list_item(self, text: str) -> str:
        """Normalize a list item by removing bullets and list numbering, but preserve quantities."""
        # Only remove obvious list numbering (digit followed by period/parenthesis and space)
        # This preserves quantities like "2 worksheets" or "16 oz jar"
        text = re.sub(r'^\d+[\.\)]\s+', '', text)  # Remove "1. " or "2) " but not "2 " or "16 "
        text = text.lstrip('•·-*◦').strip()
        return text

    def extract_careers(self, slide, debug: bool = False) -> List[str]:
        """Extract careers from a slide (usually the last slide)."""
        careers: List[str] = []
        trigger_patterns = [r"\brelated careers\b", r"\bcareers\b"]
        
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            text_frame = shape.text_frame
            full_text = text_frame.text
            if not full_text:
                continue

            if any(re.search(pat, full_text, re.IGNORECASE) for pat in trigger_patterns):
                if debug:
                    print("    Careers trigger found")
                # Line-by-line parse
                lines = [ln.strip() for ln in full_text.split('\n')]
                found_trigger = False
                for line in lines:
                    if any(re.search(pat, line, re.IGNORECASE) for pat in trigger_patterns):
                        found_trigger = True
                        continue
                    if found_trigger and line:
                        if (line.startswith(('•', '·', '-', '*', '◦')) or
                            re.match(r'^\d+\.?\s+', line) or
                            # short title-case phrases are typical career items
                            (len(line) <= 80 and line[:1].isupper() and not line.endswith(':'))):
                            clean_item = self._clean_list_item(line)
                            if clean_item and clean_item not in careers:
                                careers.append(clean_item)

                # Paragraph-by-paragraph fallback
                trigger_found = False
                for p in text_frame.paragraphs:
                    para = p.text.strip()
                    if any(re.search(pat, para, re.IGNORECASE) for pat in trigger_patterns):
                        trigger_found = True
                        continue
                    if trigger_found and para:
                        if (para.startswith(('•', '·', '-', '*', '◦')) or
                            re.match(r'^\d+\.?\s+', para) or
                            (len(para) <= 80 and para[:1].isupper() and not para.endswith(':'))):
                            clean_item = self._clean_list_item(para)
                            if clean_item and clean_item not in careers:
                                careers.append(clean_item)

        return careers

    def extract_session_materials(self, slide, debug: bool = False) -> List[str]:
        """Extract session materials from a slide. Can appear multiple times across a deck."""
        materials: List[str] = []
        trigger_patterns = [
            r"gather\s+the\s+following\s+items?",
            r"locate\s+the\s+following",
            r"\blocate\s+the\b",
        ]

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            text_frame = shape.text_frame
            full_text = text_frame.text
            if not full_text:
                continue

            if any(re.search(pat, full_text, re.IGNORECASE) for pat in trigger_patterns):
                if debug:
                    print("    Materials trigger found")
                # First pass: line-by-line following trigger
                lines = [ln.strip() for ln in full_text.split('\n')]
                found_trigger = False
                for line in lines:
                    if any(re.search(pat, line, re.IGNORECASE) for pat in trigger_patterns):
                        found_trigger = True
                        continue
                    if found_trigger and line:
                        if (line.startswith(('•', '·', '-', '*', '◦')) or
                            re.match(r'^\d+\.?\s+', line) or
                            # Typical material entries are short
                            (len(line) <= 120 and not line.endswith(':'))):
                            clean_item = self._clean_list_item(line)
                            if clean_item and clean_item not in materials:
                                materials.append(clean_item)

                # Paragraph fallback
                trigger_found = False
                for p in text_frame.paragraphs:
                    para = p.text.strip()
                    if any(re.search(pat, para, re.IGNORECASE) for pat in trigger_patterns):
                        trigger_found = True
                        continue
                    if trigger_found and para:
                        if (para.startswith(('•', '·', '-', '*', '◦')) or
                            re.match(r'^\d+\.?\s+', para) or
                            (len(para) <= 120 and not para.endswith(':'))):
                            clean_item = self._clean_list_item(para)
                            if clean_item and clean_item not in materials:
                                materials.append(clean_item)

        return materials

    def extract_all_content_from_ppt(self, ppt_path: str, debug: bool = False) -> Dict[str, any]:
        """Extract vocabulary, goals, assessments, careers, and session materials from a PowerPoint file."""
        try:
            presentation = Presentation(ppt_path)
            all_vocabulary = []
            session_goals = []
            assessments = []
            careers: List[str] = []
            materials: List[str] = []
            
            print(f"Processing {os.path.basename(ppt_path)}...")
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                # Extract vocabulary
                vocab_items = self.extract_vocabulary_from_slide(slide, presentation, debug)
                if vocab_items:
                    print(f"  Found {len(vocab_items)} vocabulary item(s) on slide {slide_num}")
                    all_vocabulary.extend(vocab_items)
                
                # Extract session goals (can appear on multiple slides)
                goals = self.extract_session_goals(slide, debug)
                if goals:
                    print(f"  Found {len(goals)} session goal(s) on slide {slide_num}")
                    # Add new goals, avoiding duplicates
                    for goal in goals:
                        if goal not in session_goals:
                            session_goals.append(goal)
                
                # Extract assessments (can appear on multiple slides)
                assessment_items = self.extract_assessments(slide, debug)
                if assessment_items:
                    print(f"  Found {len(assessment_items)} assessment item(s) on slide {slide_num}")
                    # Add new assessments, avoiding duplicates
                    for assessment in assessment_items:
                        if assessment not in assessments:
                            assessments.append(assessment)

                # Extract careers (prefer the last occurrence; usually the last slide)
                careers_items = self.extract_careers(slide, debug)
                if careers_items:
                    print(f"  Found {len(careers_items)} career(s) on slide {slide_num}")
                    careers = careers_items  # Always override to keep the latest occurrence

                # Extract session materials (can appear multiple times)
                materials_items = self.extract_session_materials(slide, debug)
                if materials_items:
                    print(f"  Found {len(materials_items)} material item(s) on slide {slide_num}")
                    materials.extend(materials_items)
            
            return {
                'vocabulary': all_vocabulary,
                'goals': session_goals,
                'assessments': assessments,
                'careers': careers,
                'materials': materials
            }
            
        except Exception as e:
            print(f"Error processing {ppt_path}: {e}")
            return {'vocabulary': [], 'goals': [], 'assessments': [], 'careers': [], 'materials': []}
    
    def extract_module_acronym(self, filename: str) -> str:
        """Extract the 4-letter acronym from filename."""
        # Extract acronym before the first underscore
        base_name = os.path.splitext(filename)[0]
        parts = base_name.split('_')
        if parts:
            return parts[0]
        return "UNKN"
    
    def process_all_ppt_files(self, directory: str, debug: bool = False) -> Dict[str, Dict[str, any]]:
        """Process all PowerPoint files in the directory."""
        ppt_files = list(Path(directory).glob("*.pptx"))
        all_content = defaultdict(lambda: {
            'vocabulary': [],
            'goals': [],
            'assessments': [],
            'careers': [],
            'materials': []
        })
        
        if not ppt_files:
            print("No PowerPoint files found in the directory.")
            return all_content
        
        print(f"Found {len(ppt_files)} PowerPoint files to process:")
        for ppt_file in ppt_files:
            print(f"  - {ppt_file.name}")
        
        print("\nProcessing files...")
        
        for ppt_file in ppt_files:
            content = self.extract_all_content_from_ppt(str(ppt_file), debug)
            if any(content.values()):  # If any content was found
                all_content[ppt_file.name] = content
        
        return all_content
    
    def create_word_document(self, all_content: Dict[str, Dict[str, any]], 
                           output_path: str, module_acronym: str):
        """Create a Word document with extracted vocabulary, goals, assessments, careers, and materials."""
        doc = Document()
        
        # Add title
        title = doc.add_heading(f'{module_acronym} - Documentation Updates & Tickets', 0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add subtitle
        subtitle = doc.add_paragraph(f'Extracted from {module_acronym} PowerPoint Sessions')
        subtitle.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        doc.add_paragraph()  # Add spacing
        
        # Collect all content
        all_vocab = []
        all_goals = []
        all_assessments = []
        all_careers = []
        all_materials = []
        
        for filename, content in all_content.items():
            for word, definition in content['vocabulary']:
                all_vocab.append((word, definition, filename))
            for goal in content['goals']:
                all_goals.append((goal, filename))
            for assessment in content['assessments']:
                all_assessments.append((assessment, filename))
            for career in content.get('careers', []):
                all_careers.append((career, filename))
            for material in content.get('materials', []):
                all_materials.append((material, filename))
        
        # Add summary
        summary_text = f"Content Summary:\n"
        summary_text += f"• Vocabulary Terms: {len(all_vocab)}\n"
        summary_text += f"• Session Goals: {len(all_goals)}\n"
        summary_text += f"• Assessment Items: {len(all_assessments)}\n"
        summary_text += f"• Careers: {len(all_careers)}\n"
        summary_text += f"• Session Materials: {len(all_materials)}\n"
        summary_text += f"• QA Tickets: (Manual entry section)\n"
        summary_text += f"• Notes: (Manual entry section)"
        
        summary = doc.add_paragraph(summary_text)
        summary.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        doc.add_page_break()
        
        # SECTION 1: VOCABULARY TERMS
        if all_vocab:
            vocab_heading = doc.add_heading('1. Vocabulary Terms', level=1)
            
            # Sort vocabulary alphabetically
            all_vocab.sort(key=lambda x: x[0].lower())
            
            for i, (word, definition, source_file) in enumerate(all_vocab, 1):
                # Add vocabulary word as heading
                word_heading = doc.add_heading(f'{i}. {word}', level=2)
                
                # Add definition
                definition_para = doc.add_paragraph(definition)
                
                # Add source file reference
                source_para = doc.add_paragraph(f'Source: {source_file}')
                source_para.runs[0].italic = True
                
                # Add spacing between terms
                if i < len(all_vocab):
                    doc.add_paragraph()
        
        # SECTION 2: SESSION GOALS
        if all_goals:
            doc.add_page_break()
            goals_heading = doc.add_heading('2. Session Goals', level=1)
            
            # Group goals by session
            goals_by_session = defaultdict(list)
            for goal, source_file in all_goals:
                goals_by_session[source_file].append(goal)
            
            for filename in sorted(goals_by_session.keys()):
                session_heading = doc.add_heading(f'{filename}', level=2)
                goals = goals_by_session[filename]
                
                for i, goal in enumerate(goals, 1):
                    goal_para = doc.add_paragraph(f'{i}. {goal}')
                
                doc.add_paragraph()  # Spacing between sessions
        
        # SECTION 3: ASSESSMENTS
        if all_assessments:
            doc.add_page_break()
            assessments_heading = doc.add_heading('3. Assessment Items', level=1)
            
            # Group assessments by session
            assessments_by_session = defaultdict(list)
            for assessment, source_file in all_assessments:
                assessments_by_session[source_file].append(assessment)
            
            for filename in sorted(assessments_by_session.keys()):
                session_heading = doc.add_heading(f'{filename}', level=2)
                assessments = assessments_by_session[filename]
                
                for i, assessment in enumerate(assessments, 1):
                    assessment_para = doc.add_paragraph(f'{i}. {assessment}')
                
                doc.add_paragraph()  # Spacing between sessions

        # SECTION 4: CAREERS (from the last slide of each session)
        if all_careers:
            doc.add_page_break()
            careers_heading = doc.add_heading('4. Related Careers', level=1)
            careers_by_session = defaultdict(list)
            for career, source_file in all_careers:
                careers_by_session[source_file].append(career)

            for filename in sorted(careers_by_session.keys()):
                session_heading = doc.add_heading(f'{filename}', level=2)
                for i, career in enumerate(careers_by_session[filename], 1):
                    doc.add_paragraph(f'{i}. {career}')
                doc.add_paragraph()

        # SECTION 5: SESSION MATERIALS (can appear multiple times per session)
        if all_materials:
            doc.add_page_break()
            materials_heading = doc.add_heading('5. Session Materials', level=1)
            materials_by_session = defaultdict(list)
            for material, source_file in all_materials:
                materials_by_session[source_file].append(material)

            for filename in sorted(materials_by_session.keys()):
                session_heading = doc.add_heading(f'{filename}', level=2)
                for i, material in enumerate(materials_by_session[filename], 1):
                    doc.add_paragraph(f'{i}. {material}')
                doc.add_paragraph()
        
        # Add QA Tickets and Notes sections
        doc.add_page_break()
        
        # QA Tickets section
        qa_heading = doc.add_heading('QA Tickets', level=1)
        
        # Add some blank lines for manual entry
        for _ in range(8):
            doc.add_paragraph()
        
        # Notes section
        notes_heading = doc.add_heading('Notes', level=1)
        
        # Add some blank lines for manual entry
        for _ in range(8):
            doc.add_paragraph()
        
        # Save the document
        doc.save(output_path)
        print(f"\nWord document saved as: {output_path}")
    
    def run(self, directory: str = None, output_filename: str = None, debug: bool = False):
        """Main execution method."""
        if directory is None:
            directory = os.getcwd()
        
        print(f"Scanning directory: {directory}")
        
        # Extract all content from PowerPoint files
        all_content = self.process_all_ppt_files(directory, debug)
        
        if not all_content:
            print("No content found in any PowerPoint files.")
            return
        
        # Determine module acronym
        ppt_files = list(Path(directory).glob("*.pptx"))
        if ppt_files:
            module_acronym = self.extract_module_acronym(ppt_files[0].name)
        else:
            module_acronym = input("Enter the 4-letter module acronym: ").strip()
        
        # Generate output filename if not provided
        if output_filename is None:
            output_filename = f"{module_acronym}_Doc Updates & Tickets.docx"
        
        output_path = os.path.join(directory, output_filename)
        
        # Create Word document
        self.create_word_document(all_content, output_path, module_acronym)
        
        # Print summary
        total_vocab = sum(len(content['vocabulary']) for content in all_content.values())
        total_goals = sum(len(content['goals']) for content in all_content.values())
        total_assessments = sum(len(content['assessments']) for content in all_content.values())
        total_careers = sum(len(content.get('careers', [])) for content in all_content.values())
        total_materials = sum(len(content.get('materials', [])) for content in all_content.values())
        
        print(f"\nSummary:")
        print(f"  - Files processed: {len(all_content)}")
        print(f"  - Total vocabulary terms extracted: {total_vocab}")
        print(f"  - Total session goals extracted: {total_goals}")
        print(f"  - Total assessment items extracted: {total_assessments}")
        print(f"  - Total careers extracted: {total_careers}")
        print(f"  - Total session materials extracted: {total_materials}")
        print(f"  - Output file: {output_filename}")
        
        # Print detailed breakdown
        print(f"\nDetailed breakdown:")
        for filename, content in all_content.items():
            vocab_count = len(content['vocabulary'])
            goals_count = len(content['goals'])
            assessments_count = len(content['assessments'])
            careers_count = len(content.get('careers', []))
            materials_count = len(content.get('materials', []))
            print(f"  {filename}: {vocab_count} vocab, {goals_count} goals, {assessments_count} assessments, {careers_count} careers, {materials_count} materials")
            
            if content['vocabulary']:
                for word, definition in content['vocabulary']:
                    print(f"    • Vocab: {word}")
            if content['goals']:
                for goal in content['goals']:
                    print(f"    • Goal: {goal[:50]}{'...' if len(goal) > 50 else ''}")
            if content['assessments']:
                for assessment in content['assessments']:
                    print(f"    • Assessment: {assessment[:50]}{'...' if len(assessment) > 50 else ''}")
            if content.get('careers'):
                for career in content['careers']:
                    print(f"    • Career: {career}")
            if content.get('materials'):
                for material in content['materials']:
                    print(f"    • Material: {material[:50]}{'...' if len(material) > 50 else ''}")


def main():
    """Main function to run the documentation extractor."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Extract vocabulary, goals, and assessments from PowerPoint files')
    parser.add_argument('--directory', '-d', 
                       help='Directory containing PowerPoint files (default: current directory)')
    parser.add_argument('--output', '-o', 
                       help='Output Word document filename (default: auto-generated)')
    parser.add_argument('--acronym', '-a',
                       help='Module acronym (default: extracted from first file)')
    parser.add_argument('--debug', action='store_true',
                       help='Enable debug output to see color and formatting details')
    
    args = parser.parse_args()
    
    extractor = DocumentationExtractor()
    
    # Use provided directory or current directory
    directory = args.directory or os.getcwd()
    
    # If acronym is provided, use it for output filename
    if args.acronym:
        output_filename = f"{args.acronym}_Doc Updates & Tickets.docx"
    else:
        output_filename = args.output
    
    extractor.run(directory, output_filename, args.debug)


if __name__ == "__main__":
    main()
